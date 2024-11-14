from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
import fitz

def ft_extract_xlsx(pathname: str) -> str:
    try:
        wb = load_workbook(pathname)
        content = ""
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            for row in ws.iter_rows(values_only=True):
                content += "\t".join(str(cell) for cell in row) + "\n"
        return content
    except Exception as e:
        print(f"Error reading Excel file: {e}")
        return ""

def ft_extract_pdf(pathname: str) -> str:
    try:
        pdf_document = fitz.open(pathname)
        content = ""
        for page_num in range(len(pdf_document)):
            page = pdf_document.load_page(page_num)
            content += page.get_text()
        pdf_document.close()
        return content
    except Exception as e:
        print(f"Error reading PDF document: {e}")
        return ""

def ft_extract_docx(pathname: str) -> str:
    try:
        doc = Document(pathname)
        content = ""
        for paragraph in doc.paragraphs:
            content += paragraph.text + "\n"
        return content
    except Exception as e:
        print(f"Error reading Word document: {e}")
        return ""

def ft_extract_pptx(pathname: str) -> str:
    try:
        prs = Presentation(pathname)
        content = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    content += shape.text + "\n"
        return content
    except Exception as e:
        print(f"Error reading PowerPoint file: {e}")
        return ""
    
def ft_extract_text(pathname: str) -> str:
    with open(pathname, "r", encoding="utf-8") as file:
        content = file.read()
    return content